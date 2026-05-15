#!/usr/bin/env python3
"""文件服务 - 支持 Word"""

from flask import Flask, jsonify
from pathlib import Path
import urllib.parse
import json

app = Flask(__name__)
UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({"status": "ok"})

@app.route('/files', methods=['GET'])
def list_files():
    files = [f.name for f in UPLOAD_DIR.iterdir() if f.is_file()]
    return jsonify({"files": files, "count": len(files)})

@app.route('/docx/<path:filename>', methods=['GET'])
def read_docx(filename):
    """读取 Word 文档"""
    from docx import Document
    
    # 解码
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    print(f"请求文件: {filename}")
    print(f"完整路径: {file_path}")
    print(f"文件存在: {file_path.exists()}")
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        return jsonify({
            'filename': filename,
            'paragraphs': len(paragraphs),
            'content': paragraphs[:20],
            'full_text': '\n'.join(paragraphs)[:3000]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/ppt/<path:filename>', methods=['GET'])
def read_ppt(filename):
    from pptx import Presentation
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
            slides.append({"num": i+1, "text": "\n".join(text)})
        
        return jsonify({
            "filename": filename,
            "total_slides": len(slides),
            "slides": slides[:20],
            "full_text": "\n".join([s["text"] for s in slides])[:5000]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("文件服务启动: http://localhost:5003")
    app.run(host='0.0.0.0', port=5003, debug=False)

@app.route('/excel/<path:filename>', methods=['GET'])
def read_excel(filename):
    """读取 Excel 文件"""
    import pandas as pd
    import urllib.parse
    
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    if not file_path.suffix.lower() in ['.xlsx', '.xls', '.csv']:
        return jsonify({'error': 'Not an Excel/CSV file'}), 400
    
    try:
        # 读取所有 sheet
        if file_path.suffix.lower() == '.csv':
            df = pd.read_csv(file_path)
            sheets_data = {"Sheet1": df.to_dict(orient='records')}
        else:
            excel_file = pd.ExcelFile(file_path)
            sheets_data = {}
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheets_data[sheet_name] = df.to_dict(orient='records')[:100]  # 限制行数
        
        # 统计信息
        total_rows = sum(len(rows) for rows in sheets_data.values())
        total_cols = 0
        for rows in sheets_data.values():
            if rows:
                total_cols = max(total_cols, len(rows[0]) if isinstance(rows[0], dict) else 0)
        
        return jsonify({
            'filename': filename,
            'type': 'excel',
            'sheets': list(sheets_data.keys()),
            'sheet_count': len(sheets_data),
            'total_rows': total_rows,
            'total_columns': total_cols,
            'data': sheets_data,
            'preview': {k: v[:5] for k, v in sheets_data.items()}  # 每 sheet 前5行预览
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/excel/<path:filename>/sheet/<sheet_name>', methods=['GET'])
def read_excel_sheet(filename, sheet_name):
    """读取 Excel 指定 Sheet"""
    import pandas as pd
    import urllib.parse
    
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    try:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        return jsonify({
            'filename': filename,
            'sheet': sheet_name,
            'rows': len(df),
            'columns': list(df.columns),
            'data': df.to_dict(orient='records')[:50]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
