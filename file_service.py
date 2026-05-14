#!/usr/bin/env python3
"""独立文件服务 - 处理文件上传下载"""

from flask import Flask, jsonify, request
from pathlib import Path
import json

app = Flask(__name__)

UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({"status": "ok"})

@app.route('/files', methods=['GET'])
def list_files():
    files = [f.name for f in UPLOAD_DIR.iterdir() if f.is_file()]
    return jsonify({"files": files, "count": len(files)})

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No filename'}), 400
    save_path = UPLOAD_DIR / file.filename
    file.save(save_path)
    return jsonify({'success': True, 'filename': file.filename})

@app.route('/read/<filename>', methods=['GET'])
def read_file(filename):
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({'error': 'Not found'}), 404
    content = file_path.read_text(encoding='utf-8')
    return jsonify({'filename': filename, 'content': content})

if __name__ == '__main__':
    print("📁 文件服务启动 on port 5003")
    app.run(host='0.0.0.0', port=5003, debug=False)

@app.route('/upload/multiple', methods=['POST'])
def upload_multiple():
    """批量上传"""
    if 'files' not in request.files:
        return jsonify({'error': 'No files'}), 400
    files = request.files.getlist('files')
    uploaded = []
    for f in files:
        if f.filename:
            save_path = UPLOAD_DIR / f.filename
            f.save(save_path)
            uploaded.append(f.filename)
    return jsonify({'success': True, 'uploaded': uploaded, 'count': len(uploaded)})

@app.route('/delete/<filename>', methods=['DELETE'])
def delete_file(filename):
    """删除文件"""
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({'error': 'Not found'}), 404
    file_path.unlink()
    return jsonify({'success': True, 'deleted': filename})

@app.route('/search', methods=['GET'])
def search_files():
    """搜索文件"""
    keyword = request.args.get('q', '')
    files = []
    for f in UPLOAD_DIR.iterdir():
        if f.is_file() and keyword.lower() in f.name.lower():
            files.append(f.name)
    return jsonify({'files': files, 'count': len(files), 'keyword': keyword})

# 增强读取功能
@app.route('/read/<filename>', methods=['GET'])
def read_file_enhanced(filename):
    """智能读取多种格式文件"""
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({'error': 'Not found'}), 404
    
    ext = file_path.suffix.lower()
    content = ""
    file_type = "unknown"
    
    try:
        if ext == '.txt':
            content = file_path.read_text(encoding='utf-8')
            file_type = "text"
        
        elif ext == '.json':
            import json
            content = json.loads(file_path.read_text())
            file_type = "json"
        
        elif ext == '.pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text() + "\n"
            file_type = "pdf"
        
        elif ext in ['.xlsx', '.xls']:
            import pandas as pd
            df = pd.read_excel(file_path)
            content = df.to_dict(orient='records')
            file_type = "excel"
        
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            content = "\n".join([p.text for p in doc.paragraphs])
            file_type = "word"
        
        elif ext in ['.jpg', '.png', '.gif']:
            from PIL import Image
            img = Image.open(file_path)
            content = {"size": img.size, "mode": img.mode, "format": img.format}
            file_type = "image"
        
        else:
            content = "不支持的文件类型"
        
        return jsonify({
            'filename': filename,
            'type': file_type,
            'content': content,
            'size': file_path.stat().st_size
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/read/ppt/<filename>', methods=['GET'])
def read_ppt(filename):
    """读取 PowerPoint 文件"""
    from pptx import Presentation
    
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({'error': 'Not found'}), 404
    
    if not filename.endswith(('.pptx', '.ppt')):
        return jsonify({'error': 'Not a PowerPoint file'}), 400
    
    prs = Presentation(file_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        slides_data.append({
            "slide_num": i + 1,
            "content": "\n".join(slide_text),
            "text_length": len("\n".join(slide_text))
        })
    
    # 生成摘要
    full_text = "\n".join([s["content"] for s in slides_data])
    
    return jsonify({
        'filename': filename,
        'type': 'powerpoint',
        'total_slides': len(slides_data),
        'slides': slides_data,
        'full_text': full_text[:5000],  # 限制长度
        'summary': f"共 {len(slides_data)} 页幻灯片",
        'has_tables': any(hasattr(shape, "table") for slide in prs.slides for shape in slide.shapes)
    })

@app.route('/ppt/summary/<filename>', methods=['GET'])
def summarize_ppt(filename):
    """生成 PPT 摘要（通过 LLM）"""
    import requests
    import json
    
    # 先读取 PPT
    from pptx import Presentation
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({'error': 'Not found'}), 404
    
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text)
    
    text = "\n".join(full_text)[:3000]  # 限制长度
    
    # 调用 LLM 生成摘要
    prompt = f"请总结以下 PowerPoint 内容的核心要点:\n\n{text}"
    try:
        resp = requests.post("http://localhost:11434/api/generate",
            json={"model": "qwen2.5:7b", "prompt": prompt, "stream": False},
            timeout=120)
        summary = resp.json().get('response', '')
        return jsonify({
            'filename': filename,
            'summary': summary,
            'total_slides': len(prs.slides),
            'text_length': len(text)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/read/ppt/<path:filename>', methods=['GET'])
def read_ppt_file(filename):
    """读取 PowerPoint 文件（支持中文路径）"""
    import urllib.parse
    from pptx import Presentation
    
    # 解码文件名
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': f'File not found: {filename}'}), 404
    
    if not filename.endswith(('.pptx', '.ppt')):
        return jsonify({'error': 'Not a PowerPoint file'}), 400
    
    try:
        prs = Presentation(file_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            slides_data.append({
                "slide_num": i + 1,
                "content": "\n".join(slide_text),
                "text_length": len("\n".join(slide_text))
            })
        
        full_text = "\n".join([s["content"] for s in slides_data])
        
        return jsonify({
            'filename': filename,
            'type': 'powerpoint',
            'total_slides': len(slides_data),
            'slides': slides_data[:10],  # 只返回前10页
            'full_text': full_text[:3000],
            'summary': f"共 {len(slides_data)} 页幻灯片"
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
