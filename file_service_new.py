#!/usr/bin/env python3
"""文件服务 - 简洁版"""

from flask import Flask, jsonify, request
from pathlib import Path
import urllib.parse
import json

app = Flask(__name__)
UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({"status": "ok"})

@app.route('/files', methods=['GET'])
def list_files():
    files = [f.name for f in UPLOAD_DIR.iterdir() if f.is_file()]
    return jsonify({"files": files, "count": len(files)})

@app.route('/ppt/<path:filename>', methods=['GET'])
def read_ppt(filename):
    """读取 PowerPoint 文件"""
    from pptx import Presentation
    
    # URL 解码
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    if not file_path.suffix.lower() in ['.pptx', '.ppt']:
        return jsonify({'error': 'Not a PowerPoint file'}), 400
    
    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
            slides.append({"num": i+1, "text": "\n".join(text)})
        
        return jsonify({
            "filename": filename,
            "total_slides": len(slides),
            "slides": slides[:20],
            "full_text": "\n".join([s["text"] for s in slides])[:5000]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("文件服务启动: http://localhost:5003")
    app.run(host='0.0.0.0', port=5003, debug=False)

@app.route('/docx/<path:filename>', methods=['GET'])
def read_docx(filename):
    """读取 Word 文档"""
    from docx import Document
    import urllib.parse
    
    filename = urllib.parse.unquote(filename)
    file_path = UPLOAD_DIR / filename
    
    if not file_path.exists():
        return jsonify({'error': 'File not found'}), 404
    
    if not file_path.suffix.lower() in ['.docx', '.doc']:
        return jsonify({'error': 'Not a Word document'}), 400
    
    try:
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 提取表格内容
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        
        full_text = "\n".join(paragraphs)
        
        return jsonify({
            'filename': filename,
            'type': 'word',
            'paragraphs': len(paragraphs),
            'content': paragraphs[:50],  # 前50段
            'full_text': full_text[:5000],
            'tables': tables[:5] if tables else [],
            'has_tables': len(tables) > 0
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
